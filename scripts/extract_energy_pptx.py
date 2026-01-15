import os
import sys
from pptx import Presentation

def extract_text(pptx_path, output_path):
    if not os.path.exists(pptx_path):
        print(f"Error: Source file not found: {pptx_path}")
        sys.exit(1)

    print(f"Opening {pptx_path}...")
    prs = Presentation(pptx_path)
    content = []

    for i, slide in enumerate(prs.slides, start=1):
        print(f"Processing slide {i}...")
        content.append(f"## Slide {i}")
        
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                content.append(shape.text.strip())
            
            # Extract from tables
            if shape.has_table:
                table_text = []
                for row in shape.table.rows:
                    row_text = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if row_text:
                        table_text.append(" | ".join(row_text))
                if table_text:
                    content.append("\n".join(table_text))
        
        content.append("\n---\n")

    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    with open(output_path, "w", encoding="utf-8") as f:
        f.write("\n\n".join(content))
    
    print(f"Successfully extracted content to {output_path}")

if __name__ == "__main__":
    src = "docs/industry-materials/用友BIP超级版-能源行业解决方案.pptx"
    dst = "docs/extracted-content/energy-content.md"
    extract_text(src, dst)
