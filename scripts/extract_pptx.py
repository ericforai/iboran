import sys
import json
from pptx import Presentation

def extract_pptx_to_json(pptx_path, output_path):
    prs = Presentation(pptx_path)
    slides_data = []

    for i, slide in enumerate(prs.slides):
        slide_title = ""
        if slide.shapes.title:
            slide_title = slide.shapes.title.text
        
        body_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text != slide_title:
                body_text.append(shape.text)
        
        slides_data.append({
            "slide_number": i + 1,
            "title": slide_title,
            "content": body_text
        })

    with open(output_path, 'w', encoding='utf-8') as f:
        json.dump(slides_data, f, ensure_ascii=False, indent=2)

if __name__ == "__main__":
    if len(sys.argv) < 3:
        print("Usage: python extract_pptx.py <input_pptx> <output_json>")
        sys.exit(1)
    
    input_file = sys.argv[1]
    output_file = sys.argv[2]
    extract_pptx_to_json(input_file, output_file)
    print(f"Successfully extracted {input_file} to {output_file}")
