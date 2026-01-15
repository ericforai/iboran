import pptx
import sys

def extract_text(pptx_path):
    prs = pptx.Presentation(pptx_path)
    text_runs = []
    for i, slide in enumerate(prs.slides):
        text_runs.append(f"--- Slide {i+1} ---")
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_runs.append(shape.text)
    return "\n".join(text_runs)

if __name__ == "__main__":
    path = "docs/product-materials/2.3 管理会计：精细实时，高效智能，管理会计助力企业创造价值.pptx"
    print(extract_text(path))
